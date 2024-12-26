from abc import ABC, abstractmethod
import pandas as pd
from typing import Dict, Any, List, Optional, Tuple, BinaryIO
import logging
from src.utils.retry_utils import DataValidator
import requests
from bs4 import BeautifulSoup
import time
import random
from urllib.parse import urljoin, urlparse
import PyPDF2
import io
import fitz  # PyMuPDF
import docx
import xml.etree.ElementTree as ET
import json
import re
import arxiv
import scholarly
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import numpy as np
from PIL import Image
import pytesseract
import tempfile
import os
import zipfile
import tabula
import pdfplumber
import openpyxl
from openpyxl.drawing.image import Image as XLImage
import csv
import yaml
import toml
import configparser
import mimetypes
import chardet
import latex2mathml.converter
import markdown
import nbformat
from nbconvert import PythonExporter
import epub
from pptx import Presentation
from xml.dom import minidom
import sqlite3
import h5py
import scipy.io
import pickle
import msgpack
import base64
from astropy.io import fits
import netCDF4 as nc
import vtk
from vtk.util.numpy_support import vtk_to_numpy
import cfgrib
from PIL.ExifTags import TAGS

logger = logging.getLogger(__name__)

class FileFormatHandler:
    """Handles different file formats and their conversions."""
    
    def __init__(self):
        self.metadata_extractors = {
            'pdf': self._extract_pdf_metadata,
            'image': self._extract_image_metadata,
            'office': self._extract_office_metadata,
            'scientific': self._extract_scientific_metadata
        }
    
    def _get_mime_type(self, content: bytes, filename: str = None) -> str:
        """Determine MIME type from content and filename."""
        if filename:
            guess = mimetypes.guess_type(filename)[0]
            if guess:
                return guess
        
        # Try to detect from content
        magic_numbers = {
            b'%PDF': 'application/pdf',
            b'PK\x03\x04': 'application/zip',
            b'\x89PNG': 'image/png',
            b'\xFF\xD8\xFF': 'image/jpeg',
            b'GIF87a': 'image/gif',
            b'GIF89a': 'image/gif',
            b'{\n': 'application/json',
            b'<?xml': 'application/xml',
        }
        
        for magic, mime in magic_numbers.items():
            if content.startswith(magic):
                return mime
        
        return 'application/octet-stream'
    
    def _detect_encoding(self, content: bytes) -> str:
        """Detect text encoding of content."""
        result = chardet.detect(content)
        return result['encoding'] or 'utf-8'
    
    def extract_from_excel(self, content: bytes) -> Dict[str, Any]:
        """Extract data from Excel files (xlsx/xls)."""
        result = {'sheets': {}, 'images': [], 'metadata': {}}
        
        try:
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            result['metadata'] = {
                'sheet_names': wb.sheetnames,
                'properties': {
                    'title': wb.properties.title,
                    'creator': wb.properties.creator,
                    'created': wb.properties.created,
                    'modified': wb.properties.modified
                }
            }
            
            # Extract each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                data = []
                for row in sheet.rows:
                    data.append([cell.value for cell in row])
                if data:
                    result['sheets'][sheet_name] = pd.DataFrame(data[1:], columns=data[0])
            
            # Extract images
            for sheet in wb:
                for image in sheet._images:
                    result['images'].append({
                        'sheet': sheet.title,
                        'image': image.ref,
                        'coordinates': (image.anchor._from.col, image.anchor._from.row)
                    })
            
        except Exception as e:
            logger.error(f"Excel extraction failed: {str(e)}")
        
        return result
    
    def extract_from_csv(self, content: bytes) -> Dict[str, Any]:
        """Extract data from CSV files with dialect detection."""
        result = {'data': None, 'dialect': None}
        
        try:
            # Detect encoding
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding)
            
            # Detect dialect
            sniffer = csv.Sniffer()
            dialect = sniffer.sniff(text_content)
            has_header = sniffer.has_header(text_content)
            
            # Read CSV
            csv_file = io.StringIO(text_content)
            if has_header:
                df = pd.read_csv(csv_file, dialect=dialect)
            else:
                df = pd.read_csv(csv_file, dialect=dialect, header=None)
            
            result['data'] = df
            result['dialect'] = {
                'delimiter': dialect.delimiter,
                'has_header': has_header
            }
            
        except Exception as e:
            logger.error(f"CSV extraction failed: {str(e)}")
        
        return result
    
    def extract_from_latex(self, content: bytes) -> Dict[str, Any]:
        """Extract data from LaTeX files."""
        result = {
            'text': '',
            'equations': [],
            'figures': [],
            'tables': [],
            'citations': []
        }
        
        try:
            text = content.decode(self._detect_encoding(content))
            
            # Extract equations
            equation_patterns = [
                r'\\\[(.*?)\\\]',  # Display math
                r'\\\((.*?)\\\)',  # Inline math
                r'\$\$(.*?)\$\$',  # Display math
                r'\$(.*?)\$'       # Inline math
            ]
            
            for pattern in equation_patterns:
                equations = re.finditer(pattern, text, re.DOTALL)
                for eq in equations:
                    try:
                        math = eq.group(1)
                        mathml = latex2mathml.converter.convert(math)
                        result['equations'].append({
                            'latex': math,
                            'mathml': mathml
                        })
                    except Exception as e:
                        logger.warning(f"Failed to convert equation: {str(e)}")
            
            # Extract figures
            figures = re.finditer(r'\\includegraphics.*?{(.*?)}', text)
            result['figures'] = [fig.group(1) for fig in figures]
            
            # Extract citations
            citations = re.finditer(r'\\cite{(.*?)}', text)
            result['citations'] = [cite.group(1) for cite in citations]
            
            # Clean text
            result['text'] = re.sub(r'\\[a-zA-Z]+{.*?}', '', text)
            
        except Exception as e:
            logger.error(f"LaTeX extraction failed: {str(e)}")
        
        return result
    
    def extract_from_jupyter(self, content: bytes) -> Dict[str, Any]:
        """Extract data from Jupyter notebooks."""
        result = {
            'cells': [],
            'metadata': {},
            'python_script': None
        }
        
        try:
            notebook = nbformat.reads(content.decode('utf-8'), as_version=4)
            result['metadata'] = notebook.metadata
            
            for cell in notebook.cells:
                cell_data = {
                    'type': cell.cell_type,
                    'source': cell.source
                }
                
                if cell.cell_type == 'code':
                    cell_data['outputs'] = []
                    for output in cell.outputs:
                        if 'data' in output:
                            cell_data['outputs'].append(output.data)
                
                result['cells'].append(cell_data)
            
            # Convert to Python script
            exporter = PythonExporter()
            python_script, _ = exporter.from_notebook_node(notebook)
            result['python_script'] = python_script
            
        except Exception as e:
            logger.error(f"Jupyter notebook extraction failed: {str(e)}")
        
        return result
    
    def extract_from_epub(self, content: bytes) -> Dict[str, Any]:
        """Extract data from EPUB files."""
        result = {
            'metadata': {},
            'chapters': [],
            'images': []
        }
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.epub', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                book = epub.read_epub(temp_file.name)
                
                # Extract metadata
                result['metadata'] = {
                    'title': book.title,
                    'authors': book.get_metadata('DC', 'creator'),
                    'language': book.get_metadata('DC', 'language'),
                    'publisher': book.get_metadata('DC', 'publisher'),
                    'rights': book.get_metadata('DC', 'rights')
                }
                
                # Extract chapters
                for item in book.get_items():
                    if item.get_type() == epub.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_content(), 'html.parser')
                        result['chapters'].append({
                            'title': soup.find('title').text if soup.find('title') else None,
                            'content': soup.get_text()
                        })
                    elif item.get_type() == epub.ITEM_IMAGE:
                        result['images'].append({
                            'id': item.id,
                            'href': item.get_name(),
                            'media_type': item.media_type,
                            'content': base64.b64encode(item.get_content()).decode('utf-8')
                        })
                
        except Exception as e:
            logger.error(f"EPUB extraction failed: {str(e)}")
        
        finally:
            if 'temp_file' in locals():
                os.unlink(temp_file.name)
        
        return result
    
    def extract_from_powerpoint(self, content: bytes) -> Dict[str, Any]:
        """Extract data from PowerPoint files."""
        result = {
            'slides': [],
            'images': [],
            'notes': []
        }
        
        try:
            prs = Presentation(io.BytesIO(content))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'number': slide_num,
                    'shapes': [],
                    'text': []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_data['text'].append(shape.text)
                    if hasattr(shape, "image"):
                        image_data = {
                            'slide': slide_num,
                            'image': base64.b64encode(shape.image.blob).decode('utf-8'),
                            'content_type': shape.image.content_type
                        }
                        result['images'].append(image_data)
                
                if slide.has_notes_slide:
                    result['notes'].append({
                        'slide': slide_num,
                        'text': slide.notes_slide.notes_text_frame.text
                    })
                
                result['slides'].append(slide_data)
            
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {str(e)}")
        
        return result
    
    def extract_from_database(self, content: bytes, db_type: str = 'sqlite') -> Dict[str, Any]:
        """Extract data from database files."""
        result = {
            'tables': {},
            'schema': {},
            'metadata': {}
        }
        
        if db_type == 'sqlite':
            try:
                with tempfile.NamedTemporaryFile(suffix='.db', delete=False) as temp_file:
                    temp_file.write(content)
                    temp_file.flush()
                    
                    conn = sqlite3.connect(temp_file.name)
                    
                    # Get list of tables
                    cursor = conn.cursor()
                    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                    tables = cursor.fetchall()
                    
                    for table in tables:
                        table_name = table[0]
                        # Get schema
                        cursor.execute(f"PRAGMA table_info({table_name});")
                        result['schema'][table_name] = cursor.fetchall()
                        
                        # Get data
                        df = pd.read_sql_query(f"SELECT * FROM {table_name}", conn)
                        result['tables'][table_name] = df
                    
                    conn.close()
                    
            except Exception as e:
                logger.error(f"SQLite extraction failed: {str(e)}")
                
            finally:
                if 'temp_file' in locals():
                    os.unlink(temp_file.name)
        
        return result
    
    def extract_from_hdf5(self, content: bytes) -> Dict[str, Any]:
        """Extract data from HDF5 files."""
        result = {
            'datasets': {},
            'attributes': {},
            'structure': {}
        }
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.h5', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                with h5py.File(temp_file.name, 'r') as f:
                    def process_group(group, path='/'):
                        for key in group.keys():
                            item = group[key]
                            item_path = f"{path}{key}"
                            
                            if isinstance(item, h5py.Dataset):
                                try:
                                    result['datasets'][item_path] = item[()]
                                    result['attributes'][item_path] = dict(item.attrs)
                                except Exception as e:
                                    logger.warning(f"Failed to read dataset {item_path}: {str(e)}")
                            
                            elif isinstance(item, h5py.Group):
                                result['structure'][item_path] = list(item.keys())
                                process_group(item, f"{item_path}/")
                    
                    process_group(f)
                    
        except Exception as e:
            logger.error(f"HDF5 extraction failed: {str(e)}")
            
        finally:
            if 'temp_file' in locals():
                os.unlink(temp_file.name)
        
        return result
    
    def extract_from_mat(self, content: bytes) -> Dict[str, Any]:
        """Extract data from MAT files using scipy instead of MATLAB."""
        from scipy.io import loadmat
        import io
        
        result = {
            'variables': {},
            'metadata': {}
        }
        
        try:
            # Create a BytesIO object from the content
            mat_file = io.BytesIO(content)
            
            # Load the MAT file using scipy
            mat_contents = loadmat(mat_file)
            
            # Extract variables, excluding special keys that start with '__'
            for key, value in mat_contents.items():
                if not key.startswith('__'):
                    result['variables'][key] = {
                        'data': value.tolist() if hasattr(value, 'tolist') else value,
                        'shape': value.shape if hasattr(value, 'shape') else None,
                        'dtype': str(value.dtype) if hasattr(value, 'dtype') else type(value).__name__
                    }
            
            # Extract metadata if available
            if '__header__' in mat_contents:
                result['metadata']['header'] = str(mat_contents['__header__'])
            if '__version__' in mat_contents:
                result['metadata']['version'] = str(mat_contents['__version__'])
            if '__globals__' in mat_contents:
                result['metadata']['globals'] = str(mat_contents['__globals__'])
            
        except Exception as e:
            logger.error(f"MAT file extraction failed: {str(e)}")
        
        return result
    
    def extract_from_pickle(self, content: bytes) -> Dict[str, Any]:
        """Extract data from Python pickle files."""
        result = {
            'data': None,
            'type_info': None
        }
        
        try:
            data = pickle.loads(content)
            result['data'] = data
            result['type_info'] = type(data).__name__
            
        except Exception as e:
            logger.error(f"Pickle extraction failed: {str(e)}")
        
        return result
    
    def extract_from_msgpack(self, content: bytes) -> Dict[str, Any]:
        """Extract data from MessagePack files."""
        result = {
            'data': None,
            'type_info': None
        }
        
        try:
            data = msgpack.unpackb(content)
            result['data'] = data
            result['type_info'] = type(data).__name__
            
        except Exception as e:
            logger.error(f"MessagePack extraction failed: {str(e)}")
        
        return result
    
    def extract_from_fits(self, content: bytes) -> Dict[str, Any]:
        """Extract data from FITS (Flexible Image Transport System) files."""
        from astropy.io import fits
        
        result = {
            'headers': [],
            'data': [],
            'metadata': {},
            'wcs_info': None
        }
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.fits', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                with fits.open(temp_file.name) as hdul:
                    # Extract WCS information if available
                    from astropy.wcs import WCS
                    try:
                        wcs = WCS(hdul[0].header)
                        result['wcs_info'] = {
                            'coordinate_system': wcs.wcs.ctype,
                            'reference_pixel': wcs.wcs.crpix.tolist(),
                            'reference_value': wcs.wcs.crval.tolist(),
                            'pixel_scale': wcs.wcs.cdelt.tolist() if wcs.wcs.cdelt is not None else None
                        }
                    except Exception as e:
                        logger.warning(f"WCS extraction failed: {str(e)}")
                    
                    # Extract data and headers from each HDU
                    for i, hdu in enumerate(hdul):
                        header_dict = dict(hdu.header)
                        result['headers'].append(header_dict)
                        
                        if hdu.data is not None:
                            result['data'].append({
                                'hdu_index': i,
                                'shape': hdu.data.shape,
                                'dtype': str(hdu.data.dtype),
                                'data': hdu.data.tolist() if hdu.data.size < 1000000 else "Data too large"
                            })
                    
                    # Extract common FITS metadata
                    result['metadata'] = {
                        'instrument': hdul[0].header.get('INSTRUME'),
                        'telescope': hdul[0].header.get('TELESCOP'),
                        'observer': hdul[0].header.get('OBSERVER'),
                        'object': hdul[0].header.get('OBJECT'),
                        'date_obs': hdul[0].header.get('DATE-OBS'),
                        'exposure_time': hdul[0].header.get('EXPTIME')
                    }
        
        except Exception as e:
            logger.error(f"FITS extraction failed: {str(e)}")
        
        finally:
            if 'temp_file' in locals():
                os.unlink(temp_file.name)
        
        return result
    
    def extract_from_netcdf(self, content: bytes) -> Dict[str, Any]:
        """Extract data from NetCDF (Network Common Data Form) files."""
        import netCDF4 as nc
        
        result = {
            'variables': {},
            'dimensions': {},
            'attributes': {},
            'metadata': {}
        }
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.nc', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                with nc.Dataset(temp_file.name, 'r') as ds:
                    # Extract dimensions
                    for dim_name, dim in ds.dimensions.items():
                        result['dimensions'][dim_name] = {
                            'size': len(dim),
                            'unlimited': dim.isunlimited()
                        }
                    
                    # Extract variables
                    for var_name, var in ds.variables.items():
                        result['variables'][var_name] = {
                            'dimensions': var.dimensions,
                            'shape': var.shape,
                            'dtype': str(var.dtype),
                            'attributes': {k: var.getncattr(k) for k in var.ncattrs()},
                            'data': var[:].tolist() if var[:].size < 1000000 else "Data too large"
                        }
                    
                    # Extract global attributes
                    result['attributes'] = {k: ds.getncattr(k) for k in ds.ncattrs()}
                    
                    # Extract common metadata
                    result['metadata'] = {
                        'conventions': ds.getncattr('Conventions') if 'Conventions' in ds.ncattrs() else None,
                        'title': ds.getncattr('title') if 'title' in ds.ncattrs() else None,
                        'institution': ds.getncattr('institution') if 'institution' in ds.ncattrs() else None,
                        'source': ds.getncattr('source') if 'source' in ds.ncattrs() else None,
                        'history': ds.getncattr('history') if 'history' in ds.ncattrs() else None
                    }
        
        except Exception as e:
            logger.error(f"NetCDF extraction failed: {str(e)}")
        
        finally:
            if 'temp_file' in locals():
                os.unlink(temp_file.name)
        
        return result
    
    def extract_from_grib(self, content: bytes) -> Dict[str, Any]:
        """Extract data from GRIB (GRIdded Binary) files."""
        import cfgrib
        
        result = {
            'messages': [],
            'metadata': {},
            'coordinates': {}
        }
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.grib', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                ds = cfgrib.open_datasets(temp_file.name)
                
                for i, dataset in enumerate(ds):
                    message = {
                        'variables': {},
                        'attributes': dict(dataset.attrs),
                        'dimensions': {dim: size for dim, size in dataset.dims.items()}
                    }
                    
                    # Extract variables
                    for var_name, var in dataset.data_vars.items():
                        message['variables'][var_name] = {
                            'attributes': dict(var.attrs),
                            'dimensions': list(var.dims),
                            'data': var.values.tolist() if var.size < 1000000 else "Data too large"
                        }
                    
                    result['messages'].append(message)
                    
                    # Extract coordinates
                    result['coordinates'] = {
                        coord_name: {
                            'values': coord.values.tolist(),
                            'attributes': dict(coord.attrs)
                        }
                        for coord_name, coord in dataset.coords.items()
                    }
                    
                    # Extract common GRIB metadata
                    result['metadata'] = {
                        'centre': dataset.attrs.get('centre'),
                        'grid_type': dataset.attrs.get('gridType'),
                        'parameter': dataset.attrs.get('paramId'),
                        'step_range': dataset.attrs.get('stepRange'),
                        'time_range': dataset.attrs.get('timeRangeIndicator')
                    }
        
        except Exception as e:
            logger.error(f"GRIB extraction failed: {str(e)}")
        
        finally:
            if 'temp_file' in locals():
                os.unlink(temp_file.name)
        
        return result
    
    def extract_from_vtk(self, content: bytes) -> Dict[str, Any]:
        """Extract data from VTK (Visualization Toolkit) files."""
        import vtk
        from vtk.util.numpy_support import vtk_to_numpy
        
        result = {
            'points': None,
            'cells': None,
            'point_data': {},
            'cell_data': {},
            'metadata': {}
        }
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.vtk', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                # Create VTK reader based on file extension
                reader = vtk.vtkDataSetReader()
                reader.SetFileName(temp_file.name)
                reader.Update()
                
                data = reader.GetOutput()
                
                # Extract points
                points = data.GetPoints()
                if points:
                    points_data = vtk_to_numpy(points.GetData())
                    result['points'] = points_data.tolist()
                
                # Extract cells
                cells = data.GetCells()
                if cells:
                    cells_data = vtk_to_numpy(cells.GetData())
                    result['cells'] = cells_data.tolist()
                
                # Extract point data
                point_data = data.GetPointData()
                for i in range(point_data.GetNumberOfArrays()):
                    array = point_data.GetArray(i)
                    name = point_data.GetArrayName(i)
                    if array and name:
                        array_data = vtk_to_numpy(array)
                        result['point_data'][name] = array_data.tolist()
                
                # Extract cell data
                cell_data = data.GetCellData()
                for i in range(cell_data.GetNumberOfArrays()):
                    array = cell_data.GetArray(i)
                    name = cell_data.GetArrayName(i)
                    if array and name:
                        array_data = vtk_to_numpy(array)
                        result['cell_data'][name] = array_data.tolist()
                
                # Extract metadata
                result['metadata'] = {
                    'number_of_points': data.GetNumberOfPoints(),
                    'number_of_cells': data.GetNumberOfCells(),
                    'data_type': data.GetDataObjectType(),
                    'bounds': data.GetBounds()
                }
        
        except Exception as e:
            logger.error(f"VTK extraction failed: {str(e)}")
        
        finally:
            if 'temp_file' in locals():
                os.unlink(temp_file.name)
        
        return result
    
    def _extract_pdf_metadata(self, metadata: Dict) -> Dict[str, Any]:
        """Extract detailed metadata from PDF files."""
        return {
            'title': metadata.get('title'),
            'author': metadata.get('author'),
            'subject': metadata.get('subject'),
            'keywords': metadata.get('keywords'),
            'creator': metadata.get('creator'),
            'producer': metadata.get('producer'),
            'creation_date': metadata.get('creation_date'),
            'modification_date': metadata.get('modification_date'),
            'trapped': metadata.get('trapped'),
            'encryption': metadata.get('encryption'),
            'page_count': metadata.get('page_count')
        }
    
    def _extract_image_metadata(self, img) -> Dict[str, Any]:
        """Extract detailed metadata from image files."""
        from PIL.ExifTags import TAGS
        
        exif_data = {}
        if hasattr(img, '_getexif') and img._getexif():
            for tag_id, value in img._getexif().items():
                tag = TAGS.get(tag_id, tag_id)
                exif_data[tag] = str(value)
        
        return {
            'format': img.format,
            'mode': img.mode,
            'size': img.size,
            'info': img.info,
            'exif': exif_data,
            'dpi': img.info.get('dpi'),
            'icc_profile': bool(img.info.get('icc_profile')),
            'animation': hasattr(img, 'n_frames') and img.n_frames > 1
        }
    
    def _extract_office_metadata(self, doc) -> Dict[str, Any]:
        """Extract detailed metadata from Office documents."""
        core_props = {
            'author': getattr(doc.core_properties, 'author', None),
            'category': getattr(doc.core_properties, 'category', None),
            'comments': getattr(doc.core_properties, 'comments', None),
            'content_status': getattr(doc.core_properties, 'content_status', None),
            'created': getattr(doc.core_properties, 'created', None),
            'identifier': getattr(doc.core_properties, 'identifier', None),
            'keywords': getattr(doc.core_properties, 'keywords', None),
            'language': getattr(doc.core_properties, 'language', None),
            'last_modified_by': getattr(doc.core_properties, 'last_modified_by', None),
            'last_printed': getattr(doc.core_properties, 'last_printed', None),
            'modified': getattr(doc.core_properties, 'modified', None),
            'revision': getattr(doc.core_properties, 'revision', None),
            'subject': getattr(doc.core_properties, 'subject', None),
            'title': getattr(doc.core_properties, 'title', None),
            'version': getattr(doc.core_properties, 'version', None)
        }
        
        return {
            'core_properties': core_props,
            'custom_properties': getattr(doc, 'custom_properties', {}),
            'content_types': getattr(doc, 'content_types', []),
            'relationships': getattr(doc, 'relationships', [])
        }
    
    def _extract_scientific_metadata(self, data) -> Dict[str, Any]:
        """Extract detailed metadata from scientific data formats."""
        return {
            'dimensions': getattr(data, 'dimensions', None),
            'variables': getattr(data, 'variables', None),
            'global_attributes': getattr(data, 'attributes', None),
            'conventions': getattr(data, 'conventions', None),
            'institution': getattr(data, 'institution', None),
            'source': getattr(data, 'source', None),
            'references': getattr(data, 'references', None),
            'history': getattr(data, 'history', None),
            'comments': getattr(data, 'comments', None),
            'coordinate_systems': getattr(data, 'coordinate_systems', None)
        }

class DataExtractor:
    """Handles various data extraction methods."""
    
    def __init__(self):
        self.chrome_options = Options()
        self.chrome_options.add_argument('--headless')
        self.chrome_options.add_argument('--disable-gpu')
        self.chrome_options.add_argument('--no-sandbox')
        self.chrome_options.add_argument('--disable-dev-shm-usage')
        self.file_format_handler = FileFormatHandler()
    
    def extract_from_pdf_advanced(self, pdf_content: bytes) -> Dict[str, Any]:
        """Extract data from PDF using multiple methods."""
        text = ""
        tables = []
        images = []
        metadata = {}
        
        # Try PyMuPDF first
        try:
            with fitz.open(stream=pdf_content, filetype="pdf") as doc:
                # Extract metadata
                metadata = doc.metadata
                
                # Extract text and images
                for page in doc:
                    text += page.get_text()
                    
                    # Extract images
                    for img in page.get_images():
                        try:
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_data = {
                                'data': base_image["image"],
                                'extension': base_image["ext"],
                                'page': page.number
                            }
                            images.append(image_data)
                        except Exception as e:
                            logger.warning(f"Failed to extract image: {str(e)}")
                
                # Extract tables using different methods
                tables.extend(self._extract_tables_from_pdf(pdf_content))
                
        except Exception as e:
            logger.error(f"PyMuPDF extraction failed: {str(e)}")
            
            # Fallback to PyPDF2
            try:
                pdf_file = io.BytesIO(pdf_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                
                metadata = pdf_reader.metadata
                
            except Exception as e2:
                logger.error(f"PyPDF2 extraction failed: {str(e2)}")
        
        return {
            'text': text,
            'tables': tables,
            'images': images,
            'metadata': metadata
        }
    
    def _extract_tables_from_pdf(self, pdf_content: bytes) -> List[pd.DataFrame]:
        """Extract tables from PDF using multiple methods."""
        tables = []
        
        # Try tabula-py first
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                temp_pdf.write(pdf_content)
                temp_pdf.flush()
                
                # Extract tables using tabula
                tables.extend(tabula.read_pdf(temp_pdf.name, pages='all'))
                
        except Exception as e:
            logger.warning(f"Tabula extraction failed: {str(e)}")
        
        # Try pdfplumber as backup
        try:
            with pdfplumber.open(io.BytesIO(pdf_content)) as pdf:
                for page in pdf.pages:
                    for table in page.extract_tables():
                        if table:
                            df = pd.DataFrame(table[1:], columns=table[0])
                            tables.append(df)
                            
        except Exception as e:
            logger.warning(f"pdfplumber extraction failed: {str(e)}")
        
        return tables
    
    def extract_from_docx(self, docx_content: bytes) -> Dict[str, Any]:
        """Extract data from DOCX files."""
        doc = docx.Document(io.BytesIO(docx_content))
        
        text = ""
        tables = []
        images = []
        
        # Extract text
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract tables
        for table in doc.tables:
            data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                data.append(row_data)
            if data:
                df = pd.DataFrame(data[1:], columns=data[0])
                tables.append(df)
        
        return {
            'text': text,
            'tables': tables,
            'images': images
        }
    
    def extract_from_image(self, image_content: bytes) -> str:
        """Extract text from images using OCR."""
        try:
            image = Image.open(io.BytesIO(image_content))
            return pytesseract.image_to_string(image)
        except Exception as e:
            logger.error(f"OCR extraction failed: {str(e)}")
            return ""
    
    def extract_from_dynamic_page(self, url: str) -> Dict[str, Any]:
        """Extract data from JavaScript-rendered pages."""
        driver = None
        try:
            driver = webdriver.Chrome(options=self.chrome_options)
            driver.get(url)
            
            # Wait for content to load
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.TAG_NAME, "body"))
            )
            
            # Extract text content
            text = driver.find_element(By.TAG_NAME, "body").text
            
            # Extract tables
            tables = []
            table_elements = driver.find_elements(By.TAG_NAME, "table")
            for table in table_elements:
                try:
                    html = table.get_attribute('outerHTML')
                    df = pd.read_html(html)[0]
                    tables.append(df)
                except Exception as e:
                    logger.warning(f"Failed to extract table: {str(e)}")
            
            return {
                'text': text,
                'tables': tables,
                'html': driver.page_source
            }
            
        except Exception as e:
            logger.error(f"Dynamic page extraction failed: {str(e)}")
            return {}
            
        finally:
            if driver:
                driver.quit()
    
    def extract_from_arxiv(self, arxiv_id: str) -> Dict[str, Any]:
        """Extract data from arXiv papers."""
        try:
            search = arxiv.Search(id_list=[arxiv_id])
            paper = next(search.results())
            
            return {
                'title': paper.title,
                'authors': [author.name for author in paper.authors],
                'abstract': paper.summary,
                'pdf_url': paper.pdf_url,
                'published': paper.published,
                'updated': paper.updated,
                'doi': paper.doi,
                'categories': paper.categories
            }
            
        except Exception as e:
            logger.error(f"arXiv extraction failed: {str(e)}")
            return {}
    
    def extract_from_google_scholar(self, title: str) -> Dict[str, Any]:
        """Extract citation data from Google Scholar."""
        try:
            search_query = scholarly.search_pubs(title)
            publication = next(search_query)
            
            return {
                'title': publication.bib.get('title'),
                'authors': publication.bib.get('author', []),
                'venue': publication.bib.get('venue'),
                'year': publication.bib.get('year'),
                'citations': publication.citedby,
                'abstract': publication.bib.get('abstract'),
                'url': publication.bib.get('url')
            }
            
        except Exception as e:
            logger.error(f"Google Scholar extraction failed: {str(e)}")
            return {}
    
    def extract_from_excel(self, excel_content: bytes) -> Dict[str, Any]:
        """Extract data from Excel files."""
        return self.file_format_handler.extract_from_excel(excel_content)
    
    def extract_from_csv(self, csv_content: bytes) -> Dict[str, Any]:
        """Extract data from CSV files."""
        return self.file_format_handler.extract_from_csv(csv_content)
    
    def extract_from_latex(self, latex_content: bytes) -> Dict[str, Any]:
        """Extract data from LaTeX files."""
        return self.file_format_handler.extract_from_latex(latex_content)
    
    def extract_from_jupyter(self, jupyter_content: bytes) -> Dict[str, Any]:
        """Extract data from Jupyter notebooks."""
        return self.file_format_handler.extract_from_jupyter(jupyter_content)
    
    def extract_from_epub(self, epub_content: bytes) -> Dict[str, Any]:
        """Extract data from EPUB files."""
        return self.file_format_handler.extract_from_epub(epub_content)
    
    def extract_from_powerpoint(self, powerpoint_content: bytes) -> Dict[str, Any]:
        """Extract data from PowerPoint files."""
        return self.file_format_handler.extract_from_powerpoint(powerpoint_content)
    
    def extract_from_database(self, db_content: bytes, db_type: str = 'sqlite') -> Dict[str, Any]:
        """Extract data from database files."""
        return self.file_format_handler.extract_from_database(db_content, db_type)
    
    def extract_from_hdf5(self, hdf5_content: bytes) -> Dict[str, Any]:
        """Extract data from HDF5 files."""
        return self.file_format_handler.extract_from_hdf5(hdf5_content)
    
    def extract_from_mat(self, mat_content: bytes) -> Dict[str, Any]:
        """Extract data from MAT files."""
        return self.file_format_handler.extract_from_mat(mat_content)
    
    def extract_from_pickle(self, pickle_content: bytes) -> Dict[str, Any]:
        """Extract data from Python pickle files."""
        return self.file_format_handler.extract_from_pickle(pickle_content)
    
    def extract_from_msgpack(self, msgpack_content: bytes) -> Dict[str, Any]:
        """Extract data from MessagePack files."""
        return self.file_format_handler.extract_from_msgpack(msgpack_content)
    
    def extract_from_fits(self, fits_content: bytes) -> Dict[str, Any]:
        """Extract data from FITS files."""
        return self.file_format_handler.extract_from_fits(fits_content)
    
    def extract_from_netcdf(self, netcdf_content: bytes) -> Dict[str, Any]:
        """Extract data from NetCDF files."""
        return self.file_format_handler.extract_from_netcdf(netcdf_content)
    
    def extract_from_grib(self, grib_content: bytes) -> Dict[str, Any]:
        """Extract data from GRIB files."""
        return self.file_format_handler.extract_from_grib(grib_content)
    
    def extract_from_vtk(self, vtk_content: bytes) -> Dict[str, Any]:
        """Extract data from VTK files."""
        return self.file_format_handler.extract_from_vtk(vtk_content)

class BaseScraper(ABC):
    """Base class for all scrapers with comprehensive data collection."""
    
    def __init__(self, output_dir: str):
        self.output_dir = output_dir
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })
        self.validator = DataValidator()
        self.extractor = DataExtractor()
        self.logger = logging.getLogger(self.__class__.__name__)
    
    def get_search_urls(self) -> List[str]:
        """Return list of URLs to scrape."""
        pass
    
    def extract_pdf_text(self, pdf_url: str) -> Optional[str]:
        """Extract full text from PDF URL."""
        try:
            response = self.session.get(pdf_url, timeout=30)
            pdf_file = io.BytesIO(response.content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PDF text from {pdf_url}: {str(e)}")
            return None
    
    def get_full_text(self, url: str, pdf_url: Optional[str] = None) -> Optional[str]:
        """Get full text from either HTML or PDF."""
        try:
            if pdf_url:
                text = self.extract_pdf_text(pdf_url)
                if text:
                    return text
            
            # Fallback to HTML if PDF fails or isn't available
            response = self.session.get(url, timeout=30)
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Try common content selectors
            content_selectors = [
                'article',
                '.paper-content',
                '.article-content',
                '#main-content',
                '.full-text'
            ]
            
            for selector in content_selectors:
                content = soup.select_one(selector)
                if content:
                    return content.get_text(strip=True, separator="\n")
            
            return None
            
        except Exception as e:
            logger.error(f"Error getting full text from {url}: {str(e)}")
            return None
    
    def get_supplementary_materials(self, url: str, soup: BeautifulSoup) -> List[Dict[str, str]]:
        """Get links to supplementary materials."""
        materials = []
        
        # Common patterns for supplementary materials
        patterns = [
            ('a[href$=".pdf"]', 'PDF'),
            ('a[href$=".doc"]', 'Document'),
            ('a[href$=".docx"]', 'Document'),
            ('a[href$=".xlsx"]', 'Spreadsheet'),
            ('a[href$=".zip"]', 'Archive'),
            ('a[href*="supplementary"]', 'Supplementary'),
            ('a[href*="supporting"]', 'Supporting'),
            ('a[href*="dataset"]', 'Dataset')
        ]
        
        for selector, material_type in patterns:
            for link in soup.select(selector):
                href = link.get('href')
                if href:
                    materials.append({
                        'url': urljoin(url, href),
                        'type': material_type,
                        'title': link.get_text(strip=True) or href.split('/')[-1]
                    })
        
        return materials
    
    def get_data_tables(self, soup: BeautifulSoup) -> List[pd.DataFrame]:
        """Extract data tables from HTML."""
        tables = []
        for table in soup.find_all('table'):
            try:
                df = pd.read_html(str(table))[0]
                tables.append(df)
            except Exception as e:
                logger.warning(f"Failed to parse table: {str(e)}")
        return tables
    
    def get_figures(self, url: str, soup: BeautifulSoup) -> List[Dict[str, str]]:
        """Extract figures and their captions."""
        figures = []
        for fig in soup.find_all(['figure', 'img']):
            try:
                img = fig.find('img') if fig.name == 'figure' else fig
                if img and img.get('src'):
                    figures.append({
                        'url': urljoin(url, img['src']),
                        'caption': fig.find('figcaption').get_text(strip=True) if fig.find('figcaption') else None,
                        'alt': img.get('alt', '')
                    })
            except Exception as e:
                logger.warning(f"Failed to parse figure: {str(e)}")
        return figures
    
    def extract_record(self, url: str) -> Dict[str, Any]:
        """Extract a complete record using multiple methods."""
        try:
            record = {}
            
            # Get initial HTML content
            response = self.session.get(url, timeout=30)
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Try to find PDF link
            pdf_url = None
            for link in soup.find_all('a', href=True):
                if link['href'].lower().endswith('.pdf'):
                    pdf_url = urljoin(url, link['href'])
                    break
            
            # Extract from PDF if available
            if pdf_url:
                pdf_response = self.session.get(pdf_url)
                pdf_data = self.extractor.extract_from_pdf_advanced(pdf_response.content)
                record.update({
                    'full_text': pdf_data['text'],
                    'pdf_tables': pdf_data['tables'],
                    'pdf_images': pdf_data['images'],
                    'pdf_metadata': pdf_data['metadata']
                })
            
            # Extract from HTML
            record.update({
                'html_text': soup.get_text(separator="\n", strip=True),
                'html_tables': self.get_data_tables(soup),
                'html_figures': self.get_figures(url, soup)
            })
            
            # Try dynamic page extraction for JavaScript content
            dynamic_data = self.extractor.extract_from_dynamic_page(url)
            if dynamic_data:
                record['dynamic_content'] = dynamic_data
            
            # Extract supplementary materials
            record['supplementary_materials'] = self.get_supplementary_materials(url, soup)
            
            # Try to get citation data
            if 'title' in record:
                scholar_data = self.extractor.extract_from_google_scholar(record['title'])
                if scholar_data:
                    record['citations'] = scholar_data
            
            # Try to get arXiv data if applicable
            arxiv_id = self._extract_arxiv_id(url)
            if arxiv_id:
                arxiv_data = self.extractor.extract_from_arxiv(arxiv_id)
                if arxiv_data:
                    record.update(arxiv_data)
            
            return record
            
        except Exception as e:
            logger.error(f"Error extracting record from {url}: {str(e)}")
            return {}
    
    def _extract_arxiv_id(self, url: str) -> Optional[str]:
        """Extract arXiv ID from URL or text."""
        arxiv_patterns = [
            r'arxiv.org/abs/(\d+\.\d+)',
            r'arxiv:(\d+\.\d+)'
        ]
        
        for pattern in arxiv_patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        
        return None
    
    def process_record(self, url: str) -> Optional[Dict[str, Any]]:
        """Process and validate a single record."""
        try:
            record = self.extract_record(url)
            
            # Validate record
            is_valid, missing_fields = self.validator.validate_record(record)
            if not is_valid:
                logger.error(f"Invalid record from {url}. Missing fields: {missing_fields}")
                return None
            
            # Enrich with optional fields
            record = self.validator.enrich_record(record)
            
            return record
            
        except Exception as e:
            logger.error(f"Error processing record from {url}: {str(e)}")
            return None
    
    def fetch_data(self) -> List[Dict[str, Any]]:
        """Fetch data from all sources with complete record collection."""
        all_records = []
        
        for url in self.get_search_urls():
            try:
                record = self.process_record(url)
                if record:
                    all_records.append(record)
                
                # Be nice to the server
                time.sleep(random.uniform(1, 3))
                
            except Exception as e:
                logger.error(f"Error fetching data from {url}: {str(e)}")
        
        return all_records
    
    def preprocess_data(self, data: List[Dict[str, Any]]) -> pd.DataFrame:
        """Convert records to DataFrame with validation."""
        if not data:
            logger.warning("No valid records found")
            return pd.DataFrame()
        
        df = pd.DataFrame(data)
        
        # Ensure all required fields are present
        missing_columns = set(self.validator.REQUIRED_FIELDS.keys()) - set(df.columns)
        if missing_columns:
            logger.error(f"Missing required columns in DataFrame: {missing_columns}")
        
        return df
